"""
datareader.py

This module contains functions to read data from a URL and process it
based on the file type. It supports various file types such as images,
audio, text, DOCX, PDF, and Excel files. The data is processed using
the OpenAI API for image analysis and audio transcription, and the
PDFMiner library for PDF text extraction.

Functions:
- datareader: Read data from a URL and process it based on the
    file type.
- process_file_type: Process the file based on its type.
- url_to_filename: Convert a URL into a safe filename.
- cache_db: Load or save data to MongoDB based on the given mode.
- handle_image: Process image data using the OpenAI API for image
    analysis.
- convert_image_to_base64: Convert image data into a base64 encoded
    string.
- handle_audio: Process audio data using the OpenAI Whisper API for
    transcription.
- handle_docx: Process DOCX files to extract text.
- handle_pdf: Process PDF files to extract text.
- handle_excel: Process Excel files to extract text.
- count_tokens: Count the number of tokens in the given text for the
    specified model.

Attributes:
    slack_bot_user_id (str): The user ID of the Slack bot.
    thread_storage (Collection): The MongoDB collection for storing
        thread information.
    url_storage (Collection): The MongoDB collection for storing URL
        information.
    aiclient (OpenAI): The OpenAI client instance.
    slack_web_client (WebClient): The Slack WebClient instance.
"""

import logging
import pandas as pd
import docx
from pptx import Presentation
from io import BytesIO
from pillow_heif import register_heif_opener

from utils.logging_utils import log_error, log_message
from envbase import (
    thread_storage,
    url_storage,
)

register_heif_opener()


def cache_db(
        url: str = None,
        thread_id: str = None,
        channel_id: str = None,
        text: str = None,
        mode: str = "save",
        url_bool: bool = False,
) -> str:
    """
    Function to either load or save data to MongoDB based
    on the given mode.

    This function handles the caching of data by either loading
    the content from MongoDB or saving the content to MongoDB.
    The collection is determined by the url_storage flag.

    Args:
        url (str): The URL of the data to be cached.
        thread_id (str): The identifier for the thread handling the
            request, used as part of the unique identifier.
        channel_id (str): The identifier for the channel.
        text (str): The text content to be saved. Defaults to None.
        mode (str): The mode of operation, either 'load' or 'save'
            ('save' by default).
        url_storage (bool): If True, use the url_storage collection.
            Defaults to False, which uses the thread_storage collection.

    Returns:
        str: The text loaded from the database if in 'load' mode,
            otherwise None.

    Raises:
        ValueError: If an invalid mode is specified.
        Exception: If an error occurs while accessing MongoDB.
    """
    try:
        if url_bool:
            # Use the url_storage collection
            collection = url_storage
            query = {"url": url}
        else:
            # Use the thread_storage collection
            collection = thread_storage
            query = {"thread_id": thread_id,
                     "channel_id": channel_id,
                     "url": url,
            }
        if mode == "save" and text is not None:
            # Save the text to the specified collection
            data = {**query, "text": text}
            collection.update_one(query, {"$set": data}, upsert=True)
            return

        elif mode == "load":
            # Load the text from the specified collection
            document = collection.find_one(query)
            if document:
                return document.get("text")
            else:
                return None

        else:
            raise ValueError("Invalid mode specified. Use 'load' or 'save'.")

    except Exception as e:
        log_error(e, "Error accessing MongoDB.")
        return None
    

async def datareader(file_data: bytes, file_type: str) -> tuple[str, str]:
    """
    Determines how to read the file based on Slack's filetype or extension.
    
    Args:
        file_data (bytes): The raw file content downloaded from Slack.
        file_type (str): The 'filetype' field from Slack (e.g., 'text', 'python', 'docx').
    
    Returns:
        tuple[str, str]: (The detected category, The extracted text content)
    """
    ft = (file_type or "").lower().replace(".", "")
    
    code_and_text = [
        "text", "txt", "md", "markdown", "log", "rst", "adoc", 
        "diff", "patch", "properties",
        "json", "xml", "yaml", "yml", "toml", "ini", "env",
        "python", "py", 
        "javascript", "js", "typescript", "ts", "jsx", "tsx",
        "java", "jar", "class", "kotlin", "kt", "scala",
        "c", "cpp", "c++", "cc", "h", "hpp", "csharp", "cs",
        "go", "golang", "rust", "rs", "ruby", "rb", "php",
        "swift", "r", "perl", "pl", "lua", "clojure", "groovy",
        "shell", "sh", "bash", "zsh", "fish", "bat", "cmd", "ps1", "powershell",
        "sql", "dockerfile", "makefile", "gradle", "pom", 
        "html", "htm", "css", "scss", "sass", "less",
    ]
    
    spreadsheet = ["xlsx", "xls", "xlsm", "odf", "csv", "tsv"]
    
    document = ["docx", "doc", "rtf"] 
    
    presentation = ["pptx", "ppt"]
    
    extracted_text = ""
    category = "text"

    try:
        if ft in code_and_text:
            category = "text"
            extracted_text = await handle_plain_text(file_data, ft)
        
        elif ft in document:
            category = "docx"
            extracted_text = await handle_docx(file_data)
        
        elif ft in spreadsheet:
            category = "spreadsheet"
            extracted_text = await handle_spreadsheet(file_data, ft)
        
        elif ft in presentation:
            category = "pptx"
            extracted_text = await handle_pptx(file_data)
            
        else:
            # Fallback: If Slack sends a weird type (e.g., 'erlang'),
            # we assume it's text and try to read it anyway.
            category = "unknown_text"
            extracted_text = await handle_plain_text(file_data, ft)

    except Exception as e:
        log_message(f"Error processing file type {ft}: {e}")
        extracted_text = f"Error: Could not process file of type {ft}."

    return category, extracted_text


async def handle_plain_text(file_data: bytes, file_type: str) -> str:
    """
    Decodes raw bytes into string. Handles encoding issues (UTF-8 vs Latin-1).
    """
    try:
        # Try standard UTF-8 first (Most modern files)
        return file_data.decode('utf-8')
    except UnicodeDecodeError:
        try:
            # Try Latin-1 (Common fallback for older Windows files, logs, or CSVs)
            return file_data.decode('latin-1')
        except Exception:
            return f"[Error: Could not decode text file of type {file_type}. It might be binary data.]"


async def handle_spreadsheet(file_data: bytes, file_type: str) -> str:
    """
    Reads Excel or CSV bytes and returns a markdown-compatible string representation.
    """
    source = BytesIO(file_data)
    
    try:
        # Check if it looks like a CSV or an Excel file
        if file_type in ['csv', 'tsv', 'txt', 'text']:
            # For CSV, we can try pandas, but sometimes it fails on bad formatting.
            try:
                df = pd.read_csv(source)
            except:
                # Fallback: Just read it as raw text if pandas fails on a CSV
                return await handle_plain_text(file_data, file_type)
        else:
            # Excel files
            df = pd.read_excel(source)
        
        # Convert DataFrame to a clean string table (index=False hides row numbers)
        return df.to_string(index=False)
        
    except Exception as e:
        return f"[Error reading spreadsheet: {str(e)}]"


async def handle_docx(file_data: bytes) -> str:
    """
    Extracts text from a .docx file using BytesIO.
    """
    try:
        source = BytesIO(file_data)
        doc = docx.Document(source)
        
        full_text = []
        for para in doc.paragraphs:
            if para.text.strip(): # Only add non-empty paragraphs
                full_text.append(para.text)
                
        return "\n\n".join(full_text)
    except Exception as e:
        return f"[Error reading DOCX: {str(e)}]"


async def handle_pptx(file_data: bytes) -> str:
    """
    Extracts text from a .pptx file using BytesIO.
    """
    try:
        source = BytesIO(file_data)
        prs = Presentation(source)
        
        text_runs = []
        
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            
            # Group text by slide
            if slide_text:
                text_runs.append("\n".join(slide_text))
                    
        return "\n\n---\n\n".join(text_runs) # Separator between slides
    except Exception as e:
        return f"[Error reading PPTX: {str(e)}]"